"""
Chunk Splitter - 350-word document splitting for Turnitin processing
Intelligently splits documents into 350-word chunks while preserving context and readability.
"""

import asyncio
import json
import logging
import os
import re
import time
import uuid
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass
from enum import Enum

import aiofiles
import redis.asyncio as redis
from sqlalchemy.orm import Session

from db.database import get_database
from db.models import DocLot, DocChunk, ChunkStatus


class SplitStrategy(Enum):
    """Document splitting strategies."""
    SIMPLE_WORD_COUNT = "simple_word_count"
    SENTENCE_BOUNDARY = "sentence_boundary"
    PARAGRAPH_BOUNDARY = "paragraph_boundary"
    SEMANTIC_BOUNDARY = "semantic_boundary"
    CITATION_AWARE = "citation_aware"


@dataclass
class SplitConfig:
    """Configuration for document splitting."""
    target_words: int = 350
    min_words: int = 300
    max_words: int = 400
    strategy: SplitStrategy = SplitStrategy.CITATION_AWARE
    preserve_citations: bool = True
    preserve_paragraphs: bool = True
    overlap_words: int = 20  # Overlap between chunks for context


@dataclass
class DocumentChunk:
    """Represents a document chunk."""
    chunk_id: str
    chunk_index: int
    content: str
    word_count: int
    start_position: int
    end_position: int
    contains_citations: bool
    preserves_context: bool
    quality_score: float


@dataclass
class SplitResult:
    """Result of document splitting operation."""
    lot_id: str
    total_chunks: int
    chunks: List[DocumentChunk]
    strategy_used: SplitStrategy
    total_words: int
    average_chunk_size: float
    split_quality_score: float
    processing_time: float


class ChunkSplitter:
    """
    Production-ready document chunk splitter for Turnitin processing.

    Features:
    - Smart 350-word chunk splitting
    - Citation preservation
    - Context-aware boundaries
    - Multiple splitting strategies
    - Quality scoring
    - Overlap management
    - Academic formatting preservation
    """

    def __init__(self, config: Optional[SplitConfig] = None):
        self.config = config or SplitConfig()
        self.logger = logging.getLogger(__name__)

        # Initialize Redis for caching
        self.redis_client = redis.from_url("redis://localhost:6379", decode_responses=True)

        # Splitting statistics
        self.stats = {
            "documents_split": 0,
            "total_chunks_created": 0,
            "average_split_quality": 0.0,
            "average_processing_time": 0.0,
            "strategy_usage": {strategy.value: 0 for strategy in SplitStrategy}
        }

        # Citation patterns for academic documents
        self.citation_patterns = [
            r'\([^)]*\d{4}[^)]*\)',  # (Author, 2023)
            r'\[[^\]]*\d+[^\]]*\]',  # [1], [Author 2023]
            r'\w+\s+\(\d{4}\)',      # Author (2023)
            r'\w+\s+et\s+al\.',     # Author et al.
            r'doi:\s*[0-9.]+/[^\s]+',  # DOI citations
            r'http[s]?://[^\s]+',    # URLs
        ]

        # Sentence boundary markers
        self.sentence_endings = ['.', '!', '?', ';']

        # Paragraph markers
        self.paragraph_markers = ['\n\n', '\r\n\r\n', '\n \n']

    async def split_document(self, file_path: str, file_type: str, document_title: str = "",
                           user_id: str = None) -> SplitResult:
        """
        Split document into 350-word chunks for Turnitin processing.

        Args:
            file_path: Path to the document to split
            file_type: The type of file (e.g., 'pdf', 'docx', 'txt')
            document_title: Title of the document
            user_id: ID of the user requesting the split

        Returns:
            SplitResult: Comprehensive splitting results
        """
        start_time = time.time()
        lot_id = str(uuid.uuid4())

        try:
            self.logger.info(f"📄 Starting document split for lot {lot_id}")

            # Extract content based on file type
            if file_type == 'pdf':
                chunks = await self._split_pdf(file_path, lot_id)
            elif file_type == 'docx':
                chunks = await self._split_docx(file_path, lot_id)
            elif file_type == 'pptx':
                chunks = await self._split_pptx(file_path, lot_id)
            elif file_type == 'xlsx':
                chunks = await self._split_xlsx(file_path, lot_id)
            elif file_type == 'txt':
                async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                    document_content = await f.read()
                chunks = await self._execute_split(document_content, SplitStrategy.PARAGRAPH_BOUNDARY, lot_id)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")

            # Analyze document characteristics
            doc_analysis = await self._analyze_document(" ".join(c.content for c in chunks))

            # Choose optimal splitting strategy
            strategy = self._choose_splitting_strategy(doc_analysis)

            # Calculate quality metrics
            quality_score = self._calculate_split_quality(chunks, doc_analysis)

            # Create lot in database
            await self._create_document_lot(lot_id, document_title, user_id, len(chunks))

            # Store chunks in database
            await self._store_chunks(lot_id, chunks)

            # Calculate statistics
            total_words = sum(chunk.word_count for chunk in chunks)
            avg_chunk_size = total_words / len(chunks) if chunks else 0
            processing_time = time.time() - start_time

            # Create result
            result = SplitResult(
                lot_id=lot_id,
                total_chunks=len(chunks),
                chunks=chunks,
                strategy_used=strategy,
                total_words=total_words,
                average_chunk_size=avg_chunk_size,
                split_quality_score=quality_score,
                processing_time=processing_time
            )

            # Update statistics
            self.stats["documents_split"] += 1
            self.stats["total_chunks_created"] += len(chunks)
            self.stats["strategy_usage"][strategy.value] += 1
            self._update_average_metrics(quality_score, processing_time)

            self.logger.info(f"✅ Document split complete: {len(chunks)} chunks created")

            return result

        except Exception as e:
            self.logger.error(f"Failed to split document for lot {lot_id}: {e}")
            raise

    async def _extract_text_from_pdf(self, file_path: str) -> str:
        import PyPDF2
        from io import BytesIO
        async with aiofiles.open(file_path, 'rb') as f:
            content = await f.read()
        pdf_reader = PyPDF2.PdfReader(BytesIO(content))
        return "\n".join(page.extract_text() for page in pdf_reader.pages)

    async def _extract_text_from_docx(self, file_path: str) -> str:
        import mammoth
        from io import BytesIO
        async with aiofiles.open(file_path, 'rb') as f:
            content = await f.read()
        text_result = mammoth.extract_raw_text(BytesIO(content))
        return text_result.value

    async def _split_pdf(self, file_path: str, lot_id: str) -> List[DocumentChunk]:
        """Split PDF into chunks of 1400 characters with 50% overlap."""
        content = await self._extract_text_from_pdf(file_path)
        # Simple character-based chunking for PDFs
        char_limit = 1400
        overlap = 700
        chunks = []
        for i in range(0, len(content), char_limit - overlap):
            chunk_text = content[i:i + char_limit]
            chunk = self._create_chunk(chunk_text.split(), len(chunks), i, lot_id)
            chunks.append(chunk)
        return chunks

    async def _split_docx(self, file_path: str, lot_id: str) -> List[DocumentChunk]:
        """Split DOCX by paragraph."""
        content = await self._extract_text_from_docx(file_path)
        return await self._split_paragraph_boundary(content, lot_id)

    async def _split_pptx(self, file_path: str, lot_id: str) -> List[DocumentChunk]:
        """Split PPTX by speaker notes per slide."""
        import pptx
        prs = pptx.Presentation(file_path)
        chunks = []
        for i, slide in enumerate(prs.slides):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    chunk = self._create_chunk(notes.split(), i, 0, lot_id)
                    chunks.append(chunk)
        return chunks

    async def _split_xlsx(self, file_path: str, lot_id: str) -> List[DocumentChunk]:
        """Split XLSX by sheet summary."""
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        chunks = []
        for i, sheet_name in enumerate(xls.sheet_names):
            df = pd.read_excel(xls, sheet_name)
            summary = df.describe().to_json()
            chunk = self._create_chunk(summary.split(), i, 0, lot_id)
            chunks.append(chunk)
        return chunks

    async def _analyze_document(self, content: str) -> Dict[str, Any]:
        """Analyze document characteristics to inform splitting strategy."""
        try:
            analysis = {
                "total_words": len(content.split()),
                "total_characters": len(content),
                "paragraph_count": len([p for p in content.split('\n\n') if p.strip()]),
                "sentence_count": len([s for s in re.split(r'[.!?]', content) if s.strip()]),
                "citation_count": 0,
                "citation_density": 0.0,
                "has_academic_structure": False,
                "avg_paragraph_length": 0.0,
                "avg_sentence_length": 0.0
            }

            # Count citations
            for pattern in self.citation_patterns:
                matches = re.findall(pattern, content)
                analysis["citation_count"] += len(matches)

            # Calculate citation density
            if analysis["total_words"] > 0:
                analysis["citation_density"] = analysis["citation_count"] / analysis["total_words"]

            # Check for academic structure
            academic_markers = [
                'introduction', 'methodology', 'results', 'discussion', 'conclusion',
                'abstract', 'literature review', 'references', 'bibliography'
            ]

            content_lower = content.lower()
            academic_marker_count = sum(1 for marker in academic_markers if marker in content_lower)
            analysis["has_academic_structure"] = academic_marker_count >= 2

            # Calculate average lengths
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            if paragraphs:
                analysis["avg_paragraph_length"] = sum(len(p.split()) for p in paragraphs) / len(paragraphs)

            sentences = [s.strip() for s in re.split(r'[.!?]', content) if s.strip()]
            if sentences:
                analysis["avg_sentence_length"] = sum(len(s.split()) for s in sentences) / len(sentences)

            return analysis

        except Exception as e:
            self.logger.error(f"Error analyzing document: {e}")
            return {"total_words": len(content.split()), "citation_count": 0}

    def _choose_splitting_strategy(self, doc_analysis: Dict[str, Any]) -> SplitStrategy:
        """Choose optimal splitting strategy based on document analysis."""
        try:
            # High citation density -> citation-aware splitting
            if doc_analysis.get("citation_density", 0) > 0.05:
                return SplitStrategy.CITATION_AWARE

            # Academic structure -> paragraph boundary splitting
            if doc_analysis.get("has_academic_structure", False):
                return SplitStrategy.PARAGRAPH_BOUNDARY

            # Long paragraphs -> sentence boundary splitting
            if doc_analysis.get("avg_paragraph_length", 0) > 100:
                return SplitStrategy.SENTENCE_BOUNDARY

            # Short paragraphs -> paragraph boundary splitting
            if doc_analysis.get("avg_paragraph_length", 0) > 50:
                return SplitStrategy.PARAGRAPH_BOUNDARY

            # Default to simple word count for other cases
            return SplitStrategy.SIMPLE_WORD_COUNT

        except Exception as e:
            self.logger.error(f"Error choosing splitting strategy: {e}")
            return SplitStrategy.SIMPLE_WORD_COUNT

    async def _execute_split(self, content: str, strategy: SplitStrategy,
                           lot_id: str) -> List[DocumentChunk]:
        """Execute the document split using the chosen strategy."""
        try:
            if strategy == SplitStrategy.CITATION_AWARE:
                return await self._split_citation_aware(content, lot_id)
            elif strategy == SplitStrategy.PARAGRAPH_BOUNDARY:
                return await self._split_paragraph_boundary(content, lot_id)
            elif strategy == SplitStrategy.SENTENCE_BOUNDARY:
                return await self._split_sentence_boundary(content, lot_id)
            elif strategy == SplitStrategy.SEMANTIC_BOUNDARY:
                return await self._split_semantic_boundary(content, lot_id)
            else:  # SIMPLE_WORD_COUNT
                return await self._split_simple_word_count(content, lot_id)

        except Exception as e:
            self.logger.error(f"Error executing split with strategy {strategy}: {e}")
            # Fallback to simple splitting
            return await self._split_simple_word_count(content, lot_id)

    async def _split_citation_aware(self, content: str, lot_id: str) -> List[DocumentChunk]:
        """Split document while preserving citation integrity."""
        try:
            chunks = []
            words = content.split()
            current_chunk_words = []
            current_position = 0
            chunk_index = 0

            i = 0
            while i < len(words):
                current_chunk_words.append(words[i])

                # Check if we're approaching target length
                if len(current_chunk_words) >= self.config.min_words:

                    # Look ahead for citation patterns
                    remaining_text = ' '.join(words[i:i+10])  # Look ahead 10 words

                    has_citation_ahead = any(
                        re.search(pattern, remaining_text)
                        for pattern in self.citation_patterns
                    )

                    # If we're at target length and no citation ahead, or at max length
                    if (len(current_chunk_words) >= self.config.target_words and not has_citation_ahead) or \
                       len(current_chunk_words) >= self.config.max_words:

                        # Create chunk
                        chunk = self._create_chunk(
                            current_chunk_words, chunk_index, current_position, lot_id
                        )
                        chunks.append(chunk)

                        # Prepare for next chunk with overlap
                        overlap_words = current_chunk_words[-self.config.overlap_words:] if len(current_chunk_words) > self.config.overlap_words else []
                        current_chunk_words = overlap_words
                        current_position = i - len(overlap_words) + 1
                        chunk_index += 1

                i += 1

            # Handle remaining words
            if current_chunk_words:
                chunk = self._create_chunk(
                    current_chunk_words, chunk_index, current_position, lot_id
                )
                chunks.append(chunk)

            return chunks

        except Exception as e:
            self.logger.error(f"Error in citation-aware splitting: {e}")
            return await self._split_simple_word_count(content, lot_id)

    async def _split_paragraph_boundary(self, content: str, lot_id: str) -> List[DocumentChunk]:
        """Split document at paragraph boundaries."""
        try:
            chunks = []
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]

            current_chunk_text = ""
            current_word_count = 0
            chunk_index = 0
            start_position = 0

            for paragraph in paragraphs:
                paragraph_words = len(paragraph.split())

                # If adding this paragraph would exceed max words, finalize current chunk
                if current_word_count + paragraph_words > self.config.max_words and current_chunk_text:

                    chunk = DocumentChunk(
                        chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                        chunk_index=chunk_index,
                        content=current_chunk_text.strip(),
                        word_count=current_word_count,
                        start_position=start_position,
                        end_position=start_position + len(current_chunk_text),
                        contains_citations=self._contains_citations(current_chunk_text),
                        preserves_context=True,  # Paragraph boundaries preserve context
                        quality_score=self._calculate_chunk_quality(current_chunk_text, True)
                    )

                    chunks.append(chunk)

                    # Start new chunk
                    current_chunk_text = paragraph
                    current_word_count = paragraph_words
                    start_position += len(current_chunk_text)
                    chunk_index += 1

                else:
                    # Add paragraph to current chunk
                    if current_chunk_text:
                        current_chunk_text += "\n\n" + paragraph
                    else:
                        current_chunk_text = paragraph
                    current_word_count += paragraph_words

                # If we've reached target size, consider finalizing
                if current_word_count >= self.config.target_words:
                    chunk = DocumentChunk(
                        chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                        chunk_index=chunk_index,
                        content=current_chunk_text.strip(),
                        word_count=current_word_count,
                        start_position=start_position,
                        end_position=start_position + len(current_chunk_text),
                        contains_citations=self._contains_citations(current_chunk_text),
                        preserves_context=True,
                        quality_score=self._calculate_chunk_quality(current_chunk_text, True)
                    )

                    chunks.append(chunk)

                    # Reset for next chunk
                    current_chunk_text = ""
                    current_word_count = 0
                    start_position += len(current_chunk_text)
                    chunk_index += 1

            # Handle remaining content
            if current_chunk_text:
                chunk = DocumentChunk(
                    chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                    chunk_index=chunk_index,
                    content=current_chunk_text.strip(),
                    word_count=current_word_count,
                    start_position=start_position,
                    end_position=start_position + len(current_chunk_text),
                    contains_citations=self._contains_citations(current_chunk_text),
                    preserves_context=True,
                    quality_score=self._calculate_chunk_quality(current_chunk_text, True)
                )
                chunks.append(chunk)

            return chunks

        except Exception as e:
            self.logger.error(f"Error in paragraph boundary splitting: {e}")
            return await self._split_simple_word_count(content, lot_id)

    async def _split_sentence_boundary(self, content: str, lot_id: str) -> List[DocumentChunk]:
        """Split document at sentence boundaries."""
        try:
            chunks = []

            # Split into sentences
            sentences = re.split(r'(?<=[.!?])\s+', content)
            sentences = [s.strip() for s in sentences if s.strip()]

            current_chunk_sentences = []
            current_word_count = 0
            chunk_index = 0

            for sentence in sentences:
                sentence_words = len(sentence.split())

                # Check if adding this sentence would exceed limits
                if current_word_count + sentence_words > self.config.max_words and current_chunk_sentences:

                    # Finalize current chunk
                    chunk_text = ' '.join(current_chunk_sentences)
                    chunk = DocumentChunk(
                        chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                        chunk_index=chunk_index,
                        content=chunk_text,
                        word_count=current_word_count,
                        start_position=0,  # Would need to calculate properly
                        end_position=len(chunk_text),
                        contains_citations=self._contains_citations(chunk_text),
                        preserves_context=True,  # Sentence boundaries preserve context
                        quality_score=self._calculate_chunk_quality(chunk_text, True)
                    )

                    chunks.append(chunk)

                    # Start new chunk
                    current_chunk_sentences = [sentence]
                    current_word_count = sentence_words
                    chunk_index += 1

                else:
                    # Add sentence to current chunk
                    current_chunk_sentences.append(sentence)
                    current_word_count += sentence_words

                # Check if we should finalize at target length
                if current_word_count >= self.config.target_words:
                    chunk_text = ' '.join(current_chunk_sentences)
                    chunk = DocumentChunk(
                        chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                        chunk_index=chunk_index,
                        content=chunk_text,
                        word_count=current_word_count,
                        start_position=0,
                        end_position=len(chunk_text),
                        contains_citations=self._contains_citations(chunk_text),
                        preserves_context=True,
                        quality_score=self._calculate_chunk_quality(chunk_text, True)
                    )

                    chunks.append(chunk)

                    # Reset for next chunk
                    current_chunk_sentences = []
                    current_word_count = 0
                    chunk_index += 1

            # Handle remaining sentences
            if current_chunk_sentences:
                chunk_text = ' '.join(current_chunk_sentences)
                chunk = DocumentChunk(
                    chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                    chunk_index=chunk_index,
                    content=chunk_text,
                    word_count=current_word_count,
                    start_position=0,
                    end_position=len(chunk_text),
                    contains_citations=self._contains_citations(chunk_text),
                    preserves_context=True,
                    quality_score=self._calculate_chunk_quality(chunk_text, True)
                )
                chunks.append(chunk)

            return chunks

        except Exception as e:
            self.logger.error(f"Error in sentence boundary splitting: {e}")
            return await self._split_simple_word_count(content, lot_id)

    async def _split_semantic_boundary(self, content: str, lot_id: str) -> List[DocumentChunk]:
        """Split document at semantic boundaries (future enhancement)."""
        # For now, fall back to paragraph boundary splitting
        # This could be enhanced with NLP models for semantic segmentation
        return await self._split_paragraph_boundary(content, lot_id)

    async def _split_simple_word_count(self, content: str, lot_id: str) -> List[DocumentChunk]:
        """Simple word-count based splitting (fallback method)."""
        try:
            chunks = []
            words = content.split()
            chunk_index = 0

            for i in range(0, len(words), self.config.target_words - self.config.overlap_words):
                chunk_words = words[i:i + self.config.target_words]
                chunk_text = ' '.join(chunk_words)

                chunk = DocumentChunk(
                    chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
                    chunk_index=chunk_index,
                    content=chunk_text,
                    word_count=len(chunk_words),
                    start_position=i,
                    end_position=i + len(chunk_words),
                    contains_citations=self._contains_citations(chunk_text),
                    preserves_context=False,  # Simple splitting may break context
                    quality_score=self._calculate_chunk_quality(chunk_text, False)
                )

                chunks.append(chunk)
                chunk_index += 1

            return chunks

        except Exception as e:
            self.logger.error(f"Error in simple word count splitting: {e}")
            raise

    def _create_chunk(self, words: List[str], chunk_index: int,
                     start_position: int, lot_id: str) -> DocumentChunk:
        """Create a document chunk from word list."""
        chunk_text = ' '.join(words)

        return DocumentChunk(
            chunk_id=f"{lot_id}_chunk_{chunk_index:04d}",
            chunk_index=chunk_index,
            content=chunk_text,
            word_count=len(words),
            start_position=start_position,
            end_position=start_position + len(words),
            contains_citations=self._contains_citations(chunk_text),
            preserves_context=True,  # Citation-aware preserves context
            quality_score=self._calculate_chunk_quality(chunk_text, True)
        )

    def _contains_citations(self, text: str) -> bool:
        """Check if text contains citations."""
        for pattern in self.citation_patterns:
            if re.search(pattern, text):
                return True
        return False

    def _calculate_chunk_quality(self, text: str, preserves_context: bool) -> float:
        """Calculate quality score for a chunk."""
        quality = 0.7  # Base quality

        # Bonus for preserving context
        if preserves_context:
            quality += 0.1

        # Bonus for proper word count range
        word_count = len(text.split())
        if self.config.min_words <= word_count <= self.config.max_words:
            quality += 0.1

        # Bonus for ending at sentence boundary
        if text.rstrip().endswith(('.', '!', '?')):
            quality += 0.05

        # Bonus for starting with capital letter (complete sentence)
        if text.strip() and text.strip()[0].isupper():
            quality += 0.05

        # Penalty for very short or very long chunks
        if word_count < self.config.min_words * 0.8:
            quality -= 0.2
        elif word_count > self.config.max_words * 1.2:
            quality -= 0.1

        return min(max(quality, 0.0), 1.0)

    def _calculate_split_quality(self, chunks: List[DocumentChunk],
                               doc_analysis: Dict[str, Any]) -> float:
        """Calculate overall quality score for the split."""
        if not chunks:
            return 0.0

        # Average chunk quality
        avg_chunk_quality = sum(chunk.quality_score for chunk in chunks) / len(chunks)

        # Word count distribution quality
        word_counts = [chunk.word_count for chunk in chunks]
        avg_words = sum(word_counts) / len(word_counts)
        word_count_variance = sum((wc - avg_words) ** 2 for wc in word_counts) / len(word_counts)
        word_count_quality = max(0, 1.0 - (word_count_variance / (self.config.target_words ** 2)))

        # Context preservation quality
        context_preserved = sum(1 for chunk in chunks if chunk.preserves_context) / len(chunks)

        # Citation preservation quality
        if doc_analysis.get("citation_count", 0) > 0:
            chunks_with_citations = sum(1 for chunk in chunks if chunk.contains_citations)
            citation_quality = chunks_with_citations / len(chunks)
        else:
            citation_quality = 1.0  # No citations to preserve

        # Combine metrics
        overall_quality = (
            avg_chunk_quality * 0.4 +
            word_count_quality * 0.3 +
            context_preserved * 0.2 +
            citation_quality * 0.1
        )

        return min(max(overall_quality, 0.0), 1.0)

    async def _create_document_lot(self, lot_id: str, title: str,
                                 user_id: str, chunk_count: int):
        """Create document lot in database."""
        try:
            with get_database() as db:
                lot = DocLot(
                    id=lot_id,
                    user_id=user_id,
                    title=title or f"Document Lot {lot_id[:8]}",
                    total_chunks=chunk_count,
                    chunks_completed=0,
                    status="processing",
                    created_at=datetime.utcnow()
                )

                db.add(lot)
                db.commit()

                self.logger.info(f"📝 Document lot created: {lot_id}")

        except Exception as e:
            self.logger.error(f"Failed to create document lot {lot_id}: {e}")
            raise

    async def _store_chunks(self, lot_id: str, chunks: List[DocumentChunk]):
        """Store chunks in database."""
        try:
            with get_database() as db:
                for chunk in chunks:
                    db_chunk = DocChunk(
                        id=chunk.chunk_id,
                        lot_id=lot_id,
                        chunk_index=chunk.chunk_index,
                        content=chunk.content,
                        word_count=chunk.word_count,
                        status=ChunkStatus.OPEN,
                        quality_score=chunk.quality_score,
                        contains_citations=chunk.contains_citations,
                        created_at=datetime.utcnow()
                    )

                    db.add(db_chunk)

                db.commit()

                self.logger.info(f"💾 Stored {len(chunks)} chunks for lot {lot_id}")

        except Exception as e:
            self.logger.error(f"Failed to store chunks for lot {lot_id}: {e}")
            raise

    def _update_average_metrics(self, quality: float, processing_time: float):
        """Update average quality and processing time metrics."""
        count = self.stats["documents_split"]

        # Update average quality
        current_avg_quality = self.stats["average_split_quality"]
        self.stats["average_split_quality"] = (
            (current_avg_quality * (count - 1) + quality) / count
        )

        # Update average processing time
        current_avg_time = self.stats["average_processing_time"]
        self.stats["average_processing_time"] = (
            (current_avg_time * (count - 1) + processing_time) / count
        )

    async def get_split_status(self, lot_id: str) -> Optional[Dict[str, Any]]:
        """Get splitting status for a document lot."""
        try:
            with get_database() as db:
                lot = db.query(DocLot).filter(DocLot.id == lot_id).first()

                if not lot:
                    return None

                chunks = db.query(DocChunk).filter(DocChunk.lot_id == lot_id).all()

                return {
                    "lot_id": lot_id,
                    "title": lot.title,
                    "total_chunks": lot.total_chunks,
                    "chunks_completed": lot.chunks_completed,
                    "status": lot.status,
                    "created_at": lot.created_at.isoformat(),
                    "chunks": [
                        {
                            "chunk_id": chunk.id,
                            "chunk_index": chunk.chunk_index,
                            "word_count": chunk.word_count,
                            "status": chunk.status.value,
                            "quality_score": chunk.quality_score
                        }
                        for chunk in chunks
                    ]
                }

        except Exception as e:
            self.logger.error(f"Failed to get split status for lot {lot_id}: {e}")
            return None

    async def get_splitter_stats(self) -> Dict[str, Any]:
        """Get comprehensive splitter statistics."""
        return {
            "stats": self.stats,
            "config": {
                "target_words": self.config.target_words,
                "min_words": self.config.min_words,
                "max_words": self.config.max_words,
                "overlap_words": self.config.overlap_words,
                "strategy": self.config.strategy.value
            },
            "timestamp": time.time()
        }

    async def close(self):
        """Close splitter and cleanup resources."""
        await self.redis_client.close()


# Global chunk splitter instance
chunk_splitter = ChunkSplitter()


# Utility functions for integration
async def split_document_into_chunks(content: str, title: str = "",
                                   user_id: str = None) -> SplitResult:
    """Split document into 350-word chunks."""
    return await chunk_splitter.split_document(content, title, user_id)


async def get_document_lot_status(lot_id: str) -> Optional[Dict[str, Any]]:
    """Get status of document lot splitting."""
    return await chunk_splitter.get_split_status(lot_id)


if __name__ == "__main__":
    # Test the chunk splitter
    async def test_splitter():
        """Test chunk splitter."""
        splitter = ChunkSplitter()

        # Create sample document
        test_document = """
        This is a sample academic document that will be split into chunks for Turnitin processing.
        The document contains multiple paragraphs and citations to test the splitting algorithm.

        Academic writing often requires careful citation management (Smith, 2023). The process
        of maintaining academic integrity while creating original content is crucial for students
        and researchers alike. Modern plagiarism detection tools like Turnitin help ensure
        that written work meets the highest standards of academic honesty.

        Furthermore, the integration of artificial intelligence detection capabilities has
        become increasingly important in recent years. As AI writing tools become more
        sophisticated, educational institutions must adapt their evaluation methods accordingly.

        The methodology section of this document demonstrates how technical writing can be
        effectively processed through automated systems. Each paragraph contributes to the
        overall argument while maintaining proper academic structure and citation practices.

        In conclusion, the effective splitting of documents into manageable chunks enables
        comprehensive plagiarism and AI detection analysis while preserving the integrity
        of the original academic work.
        """

        # Test document splitting
        result = await splitter.split_document(test_document, "Test Document", "test_user")

        print(f"Split result:")
        print(f"  Lot ID: {result.lot_id}")
        print(f"  Total chunks: {result.total_chunks}")
        print(f"  Strategy used: {result.strategy_used.value}")
        print(f"  Quality score: {result.split_quality_score:.2f}")
        print(f"  Processing time: {result.processing_time:.3f}s")

        for i, chunk in enumerate(result.chunks):
            print(f"\nChunk {i + 1}:")
            print(f"  Words: {chunk.word_count}")
            print(f"  Citations: {chunk.contains_citations}")
            print(f"  Quality: {chunk.quality_score:.2f}")
            print(f"  Content preview: {chunk.content[:100]}...")

        # Get stats
        stats = await splitter.get_splitter_stats()
        print(f"\nSplitter stats: {stats}")

        await splitter.close()

    asyncio.run(test_splitter())
